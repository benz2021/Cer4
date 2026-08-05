import streamlit as st
from PIL import Image, ImageDraw, ImageFont
import pandas as pd
from io import BytesIO
import zipfile
import re
import os
import tempfile
import hashlib
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

# --- นำเข้าไลบรารีสำหรับคลิกหาพิกัด ---
try:
    from streamlit_image_coordinates import streamlit_image_coordinates
except ImportError:
    st.error("⚠️ ไม่พบไลบรารี streamlit-image-coordinates")
    st.info("กรุณาเปิด Terminal แล้วพิมพ์: pip install streamlit-image-coordinates")
    st.stop()

# ==========================================
# 🛠️ CONSTANTS
# ==========================================
EMU_PER_PIXEL = 9525
PPT_DPI = 96  # PowerPoint ใช้ 96 DPI
FONT_CACHE = {}

# ==========================================
# 🛠️ HELPER FUNCTIONS
# ==========================================

def get_system_font_path():
    """ค้นหาฟอนต์ที่ปรับขนาดได้ในระบบ"""
    paths = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
        "/usr/share/fonts/truetype/freefont/FreeSans.ttf",
        "C:\\Windows\\Fonts\\arial.ttf",
        "C:\\Windows\\Fonts\\tahoma.ttf",
        "/System/Library/Fonts/Helvetica.ttf"
    ]
    for p in paths:
        if os.path.exists(p):
            return p
    return None

def fix_thai_text_old(text):
    """จัดตำแหน่งสระและวรรณยุกต์สำหรับฟอนต์รุ่นเก่า PUA - ครอบคลุมทุกกรณี"""
    if not isinstance(text, str):
        return str(text) if pd.notna(text) else ""
    
    tone_marks = ['\u0e48', '\u0e49', '\u0e4a', '\u0e4b', '\u0e4c']
    high_tone_marks = ['\uf713', '\uf714', '\uf715', '\uf716', '\uf717']
    
    # 1. วรรณยุกต์ + สระบน
    upper_vowels = ['\u0e31', '\u0e34', '\u0e35', '\u0e36', '\u0e37', '\u0e4d']
    for i, tone in enumerate(tone_marks):
        for vowel in upper_vowels:
            text = text.replace(vowel + tone, vowel + high_tone_marks[i])
    
    # 2. วรรณยุกต์ + ป ฝ ฟ
    tall_consonants = ['ป', 'ฝ', 'ฟ']
    left_tone_marks = ['\uf70a', '\uf70b', '\uf70c', '\uf70d', '\uf70e']
    for i, tone in enumerate(tone_marks):
        for cons in tall_consonants:
            text = text.replace(cons + tone, cons + left_tone_marks[i])
    
    # 3. สระอำ
    text = text.replace('\u0e4d\u0e32', '\u0e33')
    
    # 4. ญ และ ฐ เมื่อมีสระล่าง
    YO_YING_NO_BASE = '\uF70F'
    THO_THAN_NO_BASE = '\uF700'
    text = text.replace('ญุ', YO_YING_NO_BASE + 'ุ')
    text = text.replace('ญู', YO_YING_NO_BASE + 'ู')
    text = text.replace('ญฺ', YO_YING_NO_BASE + 'ฺ')
    text = text.replace('ฐุ', THO_THAN_NO_BASE + 'ุ')
    text = text.replace('ฐู', THO_THAN_NO_BASE + 'ู')
    text = text.replace('ฐฺ', THO_THAN_NO_BASE + 'ฺ')
    
    # 5. วรรณยุกต์ + ญ
    for i, tone in enumerate(tone_marks):
        text = text.replace('ญ' + tone, YO_YING_NO_BASE + high_tone_marks[i])
    
    # 6. เกี๊ยะ, เกี๋ยว, เกื๊อ, เกื๋อ
    text = text.replace('เกี๊ย', 'เกี\uf714ย')
    text = text.replace('เกี๋ย', 'เกี\uf712ย')
    text = text.replace('เกื๊อ', 'เกื\uf714อ')
    text = text.replace('เกื๋อ', 'เกื\uf712อ')
    
    # 7. ำ้
    text = text.replace('ำ้', '\u0e33\u0e49')
    
    # 8. ณ์, ร์
    text = text.replace('ณ์', '\u0e13\u0e4c')
    text = text.replace('ร์', '\u0e23\u0e4c')
    
    # 9. ปิ๊, ปี่, ฝี่, ฟี่
    text = text.replace('ปิ๊', '\u0e1b\u0e34\u0e4a')
    text = text.replace('ปี่', '\u0e1b\u0e35\u0e48')
    text = text.replace('ฝี่', '\u0e1d\u0e35\u0e48')
    text = text.replace('ฟี่', '\u0e1f\u0e35\u0e48')
    
    return text

def fix_thai_text_new(text):
    """สำหรับฟอนต์รุ่นใหม่ (OpenType)"""
    if not isinstance(text, str):
        return str(text) if pd.notna(text) else ""
    return text

def fix_thai_text(text, font_version="ใหม่"):
    if font_version == "เก่า":
        return fix_thai_text_old(text)
    return fix_thai_text_new(text)

@st.cache_data
def cache_font_data(font_name, font_bytes):
    """Cache font data เพื่อลดการใช้ RAM"""
    return font_bytes

def get_font(font_name, size):
    """โหลดฟอนต์พร้อม Cache"""
    cache_key = f"{font_name}_{size}"
    if cache_key in FONT_CACHE:
        return FONT_CACHE[cache_key]
    
    try:
        if font_name and font_name in st.session_state.fonts_dict:
            font_data = st.session_state.fonts_dict[font_name]['data']
            font = ImageFont.truetype(BytesIO(font_data), size)
            FONT_CACHE[cache_key] = font
            return font
    except Exception:
        pass
    
    sys_path = get_system_font_path()
    if sys_path:
        try:
            font = ImageFont.truetype(sys_path, size)
            FONT_CACHE[cache_key] = font
            return font
        except:
            font = ImageFont.load_default()
            FONT_CACHE[cache_key] = font
            return font
    font = ImageFont.load_default()
    FONT_CACHE[cache_key] = font
    return font

def get_text_bbox(text, font):
    """คำนวณ bbox ของข้อความ รองรับ Pillow ทุกเวอร์ชัน"""
    try:
        # ลองใช้ getbbox (Pillow 10+)
        bbox = font.getbbox(text)
        return bbox
    except AttributeError:
        try:
            # ลองใช้ getsize (Pillow เก่า)
            width, height = font.getsize(text)
            return (0, 0, width, height)
        except:
            # สุดท้ายใช้ textlength
            draw = ImageDraw.Draw(Image.new('RGB', (1, 1)))
            width = draw.textlength(text, font=font)
            return (0, 0, int(width), int(width * 0.6))

def pixels_to_points(pixels, dpi=96):
    """แปลง pixel เป็น point"""
    return pixels * 72 / dpi

def sanitize_filename(name):
    return re.sub(r'[<>:"/\\|?*]', '_', str(name)).strip() or "certificate"

@st.cache_data
def render_certificate_cached(template_bytes, texts_tuple, row_data_tuple):
    """Cache การ render เพื่อเพิ่มประสิทธิภาพ"""
    # ต้องแปลง bytes กลับเป็น Image
    template_img = Image.open(BytesIO(template_bytes))
    # แปลง texts จาก tuple กลับเป็น list
    texts = list(texts_tuple)
    row_data = dict(row_data_tuple) if row_data_tuple else None
    
    img = template_img.copy()
    if img.mode != 'RGB':
        img = img.convert('RGB')
    draw = ImageDraw.Draw(img)
    
    for txt in texts:
        if txt['type'] == 'static':
            content = txt['text']
        else:
            if row_data and txt['column'] in row_data:
                val = row_data[txt['column']]
                content = str(val) if pd.notna(val) else ""
            else:
                content = "ตัวอย่างข้อมูล"
        
        if not content: 
            continue
        
        font_version = txt.get('font_version', 'ใหม่')
        content = fix_thai_text(content, font_version)
        font = get_font(txt.get('font_name'), txt['size'])
        
        # ใช้ bbox คำนวณตำแหน่ง
        bbox = get_text_bbox(content, font)
        text_width = bbox[2] - bbox[0]
        text_height = bbox[3] - bbox[1]
        
        # คำนวณตำแหน่งให้กึ่งกลาง
        start_x = txt['x'] - (text_width / 2)
        start_y = txt['y'] - (text_height / 2)
        
        draw.text((start_x, start_y), content, fill=txt['color'], font=font)
    
    return img

def render_certificate(template_img, texts, row_data=None):
    """Wrapper สำหรับ render_certificate_cached"""
    # เก็บ template เป็น bytes สำหรับ cache
    template_bytes = BytesIO()
    template_img.save(template_bytes, format='PNG')
    template_bytes = template_bytes.getvalue()
    
    # แปลง texts เป็น tuple สำหรับ cache
    texts_tuple = tuple(texts)
    
    # แปลง row_data เป็น tuple
    row_data_tuple = tuple(row_data.items()) if row_data else None
    
    return render_certificate_cached(template_bytes, texts_tuple, row_data_tuple)

def hex_to_rgb(hex_color):
    """แปลงสี Hex เป็น RGB tuple"""
    hex_color = hex_color.lstrip('#')
    if len(hex_color) == 6:
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    return (0, 0, 0)

def create_pptx_with_editable_text(template_img, texts, data_df):
    """สร้าง PowerPoint ที่ข้อความแก้ไขได้ - ปรับปรุงการจัดตำแหน่ง"""
    prs = Presentation()
    
    img_width, img_height = template_img.size
    prs.slide_width = img_width * EMU_PER_PIXEL
    prs.slide_height = img_height * EMU_PER_PIXEL
    
    blank_slide_layout = prs.slide_layouts[6]
    
    # Cache พื้นหลัง (สร้างครั้งเดียว)
    bg_io = BytesIO()
    template_img.save(bg_io, format="PNG")
    bg_bytes = bg_io.getvalue()
    
    for idx, row in data_df.iterrows():
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # แปะภาพพื้นหลัง (ใช้ cache)
        bg_io = BytesIO(bg_bytes)
        slide.shapes.add_picture(bg_io, 0, 0, width=prs.slide_width, height=prs.slide_height)
        
        # เพิ่มข้อความแต่ละรายการ
        for txt in texts:
            if txt['type'] == 'static':
                content = txt['text']
            else:
                if txt['column'] in row:
                    val = row[txt['column']]
                    content = str(val) if pd.notna(val) else ""
                else:
                    continue
            
            if not content: 
                continue
            
            ppt_font_version = txt.get('ppt_font_version', 'ใหม่')
            ppt_content = fix_thai_text(content, ppt_font_version)
            
            # คำนวณขนาดฟอนต์เป็น point
            font_size_pt = pixels_to_points(txt['size'])
            
            # คำนวณขนาดข้อความจริง
            font = get_font(txt.get('font_name'), txt['size'])
            bbox = get_text_bbox(content, font)
            text_width_px = bbox[2] - bbox[0]
            text_height_px = bbox[3] - bbox[1]
            
            # เพิ่ม margin
            margin_px = 20
            box_width_px = text_width_px + margin_px * 2
            box_height_px = text_height_px + margin_px * 2
            
            # แปลงเป็น EMUs
            box_width_emu = int(box_width_px * EMU_PER_PIXEL)
            box_height_emu = int(box_height_px * EMU_PER_PIXEL)
            
            center_x_emu = txt['x'] * EMU_PER_PIXEL
            center_y_emu = txt['y'] * EMU_PER_PIXEL
            
            left_emu = center_x_emu - (box_width_emu / 2)
            top_emu = center_y_emu - (box_height_emu / 2)
            
            # ป้องกันกล่องออกนอกสไลด์
            left_emu = max(0, left_emu)
            top_emu = max(0, top_emu)
            
            # สร้างกล่องข้อความ
            txBox = slide.shapes.add_textbox(
                int(left_emu), 
                int(top_emu), 
                int(box_width_emu), 
                int(box_height_emu)
            )
            
            tf = txBox.text_frame
            tf.text = ppt_content
            tf.word_wrap = False
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # จัดกึ่งกลาง
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.space_before = Pt(0)
            p.space_after = Pt(0)
            
            run = p.runs[0]
            run.font.size = Pt(font_size_pt)
            if txt.get('font_name'):
                run.font.name = txt['font_name']
            
            rgb = hex_to_rgb(txt['color'])
            run.font.color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
    
    pptx_io = BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io

# ==========================================
# 🎨 UI - STREAMLIT APP
# ==========================================
st.set_page_config(page_title="Auto Cert Pro", layout="wide")

# ตั้งค่า Session State
if "click_x" not in st.session_state: 
    st.session_state.click_x = 0
if "click_y" not in st.session_state: 
    st.session_state.click_y = 0
if 'texts' not in st.session_state: 
    st.session_state.texts = []
if 'fonts_dict' not in st.session_state: 
    st.session_state.fonts_dict = {}
if 'font_names' not in st.session_state: 
    st.session_state.font_names = []
if 'template' not in st.session_state:
    st.session_state.template = None
if 'data' not in st.session_state:
    st.session_state.data = None
if 'preview_row' not in st.session_state:
    st.session_state.preview_row = 0
if 'edit_index' not in st.session_state:
    st.session_state.edit_index = None

st.title("📜 Auto Certificate Generator")

# --- SIDEBAR ---
with st.sidebar:
    st.header("1️⃣ อัปโหลดไฟล์")
    template_file = st.file_uploader("🖼️ พื้นหลังเกียรติบัตร (JPG/PNG)", type=['jpg', 'jpeg', 'png'])
    if template_file:
        st.session_state.template = Image.open(template_file)
        st.success("✅ โหลดพื้นหลังสำเร็จ")

    st.markdown("---")
    st.header("2️⃣ จัดการฟอนต์")
    
    st.subheader("ฟอนต์สำหรับ PNG/PDF")
    font_type_png = st.radio(
        "ประเภทฟอนต์ PNG/PDF",
        ["ฟอนต์รุ่นใหม่ (OpenType)", "ฟอนต์รุ่นเก่า (PUA)"],
        key="font_type_png"
    )
    
    st.subheader("ฟอนต์สำหรับ PowerPoint")
    font_type_ppt = st.radio(
        "ประเภทฟอนต์ PowerPoint",
        ["ฟอนต์รุ่นใหม่ (OpenType)", "ฟอนต์รุ่นเก่า (PUA)"],
        key="font_type_ppt"
    )
    
    uploaded_font = st.file_uploader("🔤 อัปโหลดฟอนต์ .ttf", type=['ttf'])
    
    if uploaded_font:
        f_name = uploaded_font.name.split('.')[0]
        if f_name not in st.session_state.fonts_dict:
            st.session_state.fonts_dict[f_name] = {
                'data': uploaded_font.getvalue()
            }
            st.session_state.font_names.append(f_name)
            st.success(f"✅ เพิ่มฟอนต์ '{f_name}' แล้ว")
    
    if st.session_state.font_names:
        st.markdown("---")
        st.write("**📋 ฟอนต์ที่มี:**")
        for f in st.session_state.font_names:
            st.write(f"- {f}")

    st.markdown("---")
    st.header("3️⃣ รายชื่อข้อมูล")
    data_file = st.file_uploader("📊 ไฟล์ Excel/CSV", type=['xlsx', 'xls', 'csv'])
    if data_file:
        try:
            if data_file.name.endswith('.csv'):
                st.session_state.data = pd.read_csv(data_file)
            else:
                st.session_state.data = pd.read_excel(data_file)
            st.success(f"✅ โหลดข้อมูล {len(st.session_state.data)} รายการ")
        except Exception as e:
            st.error(f"❌ โหลดข้อมูลไม่สำเร็จ: {e}")

# ตรวจสอบ template
if st.session_state.template is None:
    st.info("👈 กรุณาอัปโหลด 'พื้นหลังเกียรติบัตร' ทางด้านซ้ายเพื่อเริ่มต้น")
    st.stop()

# --- MAIN AREA ---
st.header("📍 กำหนดตำแหน่งและข้อความ")

col_img, col_form = st.columns([1.5, 1])

with col_img:
    # แสดงตัวเลือกแถวตัวอย่าง
    if st.session_state.data is not None and not st.session_state.data.empty:
        st.session_state.preview_row = st.number_input(
            "ดูตัวอย่างจากแถวที่:", 
            0, 
            len(st.session_state.data)-1, 
            st.session_state.preview_row,
            key="preview_row_input"
        )
        preview_row = st.session_state.data.iloc[st.session_state.preview_row].to_dict()
    else:
        preview_row = None
    
    # สร้างรูปพรีวิว (ใช้ Cache)
    current_preview = render_certificate(st.session_state.template, st.session_state.texts, preview_row)
    
    st.markdown("**🖱️ คลิกที่รูปเพื่อกำหนดตำแหน่ง (ข้อความจะอยู่กึ่งกลางจุดคลิก)**")
    original_w, original_h = current_preview.size
    display_w = 700 
    ratio = original_w / display_w if original_w > display_w else 1.0
    display_img = current_preview.resize((display_w, int(original_h / ratio))) if original_w > display_w else current_preview
    
    # ใช้ streamlit_image_coordinates
    try:
        coords = streamlit_image_coordinates(display_img, key="coords")
        if coords:
            new_x = int(coords['x'] * ratio)
            new_y = int(coords['y'] * ratio)
            if new_x != st.session_state.click_x or new_y != st.session_state.click_y:
                st.session_state.click_x = new_x
                st.session_state.click_y = new_y
                st.rerun()
    except Exception as e:
        st.warning("⚠️ ระบบคลิกไม่ทำงาน กรุณาใช้ slider แทน")
        col1, col2 = st.columns(2)
        with col1:
            new_x = st.slider("X", 0, original_w, st.session_state.click_x)
            if new_x != st.session_state.click_x:
                st.session_state.click_x = new_x
        with col2:
            new_y = st.slider("Y", 0, original_h, st.session_state.click_y)
            if new_y != st.session_state.click_y:
                st.session_state.click_y = new_y
    
    st.info(f"📍 พิกัดปัจจุบัน: X={st.session_state.click_x}, Y={st.session_state.click_y}")

with col_form:
    edit_mode = st.session_state.edit_index is not None
    
    if edit_mode:
        st.subheader("✏️ แก้ไขข้อความ")
        # โหลดข้อมูลเดิมมาแสดง
        edit_txt = st.session_state.texts[st.session_state.edit_index]
        # กำหนดค่าเริ่มต้น
        default_text = edit_txt.get('text', '')
        default_column = edit_txt.get('column', None)
        default_type = 'excel' if edit_txt['type'] == 'excel' else 'พิมพ์เอง'
    else:
        st.subheader("➕ เพิ่มข้อความ")
        default_text = ''
        default_column = None
        default_type = 'พิมพ์เอง'
    
    # ตัวเลือกคอลัมน์
    cols_options = st.session_state.data.columns.tolist() if st.session_state.data is not None else []
    
    # Radio - ใช้ค่าเริ่มต้นจาก edit mode
    t_type = st.radio(
        "ชนิดข้อความ",
        ["ดึงจากไฟล์รายชื่อ", "พิมพ์เอง"],
        horizontal=True,
        key="text_type",
        index=0 if default_type == 'excel' else 1
    )
    
    if t_type == "ดึงจากไฟล์รายชื่อ":
        if not cols_options:
            st.warning("⚠️ กรุณาอัปโหลดไฟล์รายชื่อก่อน")
            selected_column = None
        else:
            # หา index ของคอลัมน์เดิม (ถ้ามี)
            default_index = 0
            if default_column and default_column in cols_options:
                default_index = cols_options.index(default_column)
            selected_column = st.selectbox(
                "เลือกหัวข้อ (คอลัมน์)",
                cols_options,
                index=default_index,
                key="excel_column_select"
            )
        t_val = ""
    else:
        # ถ้าเป็นโหมดแก้ไขและพิมพ์เอง ให้โหลดข้อความเดิม
        if edit_mode and default_type == 'พิมพ์เอง':
            t_val = st.text_input("ข้อความที่ต้องการพิมพ์", value=default_text, key="text_input")
        else:
            t_val = st.text_input("ข้อความที่ต้องการพิมพ์", key="text_input")
        selected_column = None
    
    with st.form("add_text_form", clear_on_submit=not edit_mode):
        c1, c2 = st.columns(2)
        
        if edit_mode:
            x_pos = c1.number_input("ตำแหน่ง X", value=edit_txt['x'])
            y_pos = c2.number_input("ตำแหน่ง Y", value=edit_txt['y'])
            f_size = st.slider("ขนาดฟอนต์", 10, 500, value=edit_txt['size'])
            f_color = st.color_picker("เลือกสีข้อความ", value=edit_txt['color'])
            default_font_index = 0
            if edit_txt['font_name'] in st.session_state.font_names:
                default_font_index = st.session_state.font_names.index(edit_txt['font_name'])
            selected_font = st.selectbox("เลือกฟอนต์", st.session_state.font_names, index=default_font_index)
        else:
            x_pos = c1.number_input("ตำแหน่ง X", value=st.session_state.click_x)
            y_pos = c2.number_input("ตำแหน่ง Y", value=st.session_state.click_y)
            f_size = st.slider("ขนาดฟอนต์", 10, 500, value=60)
            f_color = st.color_picker("เลือกสีข้อความ", value="#000000")
            
            if not st.session_state.font_names:
                st.warning("⚠️ กรุณาอัปโหลดฟอนต์ก่อน")
                selected_font = None
            else:
                selected_font = st.selectbox("เลือกฟอนต์", st.session_state.font_names)
        
        if edit_mode:
            submit_label = "💾 อัปเดตข้อความ"
        else:
            submit_label = "➕ เพิ่มข้อความ"
        
        submit = st.form_submit_button(submit_label)
        if submit:
            if not selected_font:
                st.error("❌ กรุณาเลือกฟอนต์")
            elif t_type == "พิมพ์เอง" and not t_val and not edit_mode:
                st.error("❌ กรุณาพิมพ์ข้อความ")
            elif t_type == "ดึงจากไฟล์รายชื่อ" and not selected_column and not edit_mode:
                st.error("❌ กรุณาเลือกคอลัมน์")
            else:
                font_version_png = "เก่า" if "เก่า" in font_type_png else "ใหม่"
                font_version_ppt = "เก่า" if "เก่า" in font_type_ppt else "ใหม่"
                
                if edit_mode and st.session_state.edit_index is not None:
                    # อัปเดตข้อความ
                    st.session_state.texts[st.session_state.edit_index] = {
                        'type': 'excel' if t_type == "ดึงจากไฟล์รายชื่อ" else 'static',
                        'text': t_val if t_type == "พิมพ์เอง" else edit_txt.get('text', ''),
                        'column': selected_column if t_type == "ดึงจากไฟล์รายชื่อ" else edit_txt.get('column', None),
                        'x': x_pos,
                        'y': y_pos,
                        'size': f_size,
                        'color': f_color,
                        'font_name': selected_font,
                        'font_version': font_version_png,
                        'ppt_font_version': font_version_ppt
                    }
                    st.session_state.edit_index = None
                    st.success("✅ อัปเดตข้อความสำเร็จ!")
                else:
                    # เพิ่มข้อความใหม่
                    st.session_state.texts.append({
                        'type': 'excel' if t_type == "ดึงจากไฟล์รายชื่อ" else 'static',
                        'text': t_val,
                        'column': selected_column,
                        'x': x_pos,
                        'y': y_pos,
                        'size': f_size,
                        'color': f_color,
                        'font_name': selected_font,
                        'font_version': font_version_png,
                        'ppt_font_version': font_version_ppt
                    })
                    st.success("✅ เพิ่มข้อความสำเร็จ!")
                st.rerun()

    # ปุ่มยกเลิกแก้ไข
    if edit_mode:
        if st.button("❌ ยกเลิกการแก้ไข"):
            st.session_state.edit_index = None
            st.rerun()

    # แสดงรายการข้อความ
    if st.session_state.texts:
        st.markdown("---")
        st.write("**📋 รายการข้อความ:**")
        for i, t in enumerate(st.session_state.texts):
            lbl = t['text'] if t['type'] == 'static' else f"📊 {t['column']}"
            col_act1, col_act2, col_act3 = st.columns([3, 1, 1])
            font_version_display = t.get('font_version', 'ใหม่')
            col_act1.write(f"{i+1}. {lbl} | ฟอนต์: {t['font_name']} (รุ่น{font_version_display}) | ขนาด: {t['size']}")
            
            if col_act2.button("✏️", key=f"edit_{i}"):
                st.session_state.edit_index = i
                st.rerun()
            
            if col_act3.button("🗑️", key=f"del_{i}"):
                st.session_state.texts.pop(i)
                if st.session_state.edit_index == i:
                    st.session_state.edit_index = None
                st.rerun()

# --- Export ---
if st.session_state.data is not None and not st.session_state.data.empty and st.session_state.texts:
    st.markdown("---")
    st.header("📦 สร้างและดาวน์โหลด")
    
    c1, c2 = st.columns(2)
    filename_col = c1.selectbox("เลือกคอลัมน์สำหรับชื่อไฟล์", st.session_state.data.columns)
    file_format = c2.radio("รูปแบบไฟล์", ["PNG", "PDF", "PowerPoint"], horizontal=True)
    
    if st.button("🚀 เริ่มสร้างทั้งหมด", type="primary"):
        with st.spinner("กำลังประมวลผล..."):
            if file_format == "PowerPoint":
                pptx_io = create_pptx_with_editable_text(
                    st.session_state.template,
                    st.session_state.texts,
                    st.session_state.data
                )
                st.success("✅ สร้างไฟล์ PowerPoint เรียบร้อย! (ข้อความแก้ไขได้)")
                st.download_button(
                    "📥 ดาวน์โหลด PowerPoint",
                    pptx_io.getvalue(),
                    "certificates.pptx",
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
            else:
                zip_buffer = BytesIO()
                with zipfile.ZipFile(zip_buffer, 'w') as zf:
                    for _, row in st.session_state.data.iterrows():
                        final_img = render_certificate(st.session_state.template, st.session_state.texts, row.to_dict())
                        img_io = BytesIO()
                        if file_format == "PNG":
                            final_img.save(img_io, format="PNG")
                            ext = "png"
                        else:
                            # ใช้ PDF คุณภาพสูง
                            final_img.save(img_io, format="PDF", resolution=300.0)
                            ext = "pdf"
                        zf.writestr(f"{sanitize_filename(row[filename_col])}.{ext}", img_io.getvalue())
                st.success("✅ สร้างไฟล์ ZIP เรียบร้อย!")
                st.download_button(
                    f"📥 ดาวน์โหลด ZIP ({file_format})",
                    zip_buffer.getvalue(),
                    "certificates.zip",
                    "application/zip"
                )

# --- คำแนะนำ ---
with st.sidebar:
    st.markdown("---")
    st.markdown("""
    ### 📖 วิธีใช้
    1. อัปโหลดพื้นหลังเกียรติบัตร
    2. เลือกประเภทฟอนต์
    3. อัปโหลดฟอนต์ .ttf
    4. อัปโหลดไฟล์ Excel/CSV
    5. คลิกบนรูปเพื่อกำหนดพิกัด
    6. เพิ่ม/แก้ไขข้อความ
    7. สร้างและดาวน์โหลด
    """)
